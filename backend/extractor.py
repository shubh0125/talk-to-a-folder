import io
import fitz  # PyMuPDF
import docx
import pptx
from models import FileInfo


def extract_text(file: FileInfo, content: bytes) -> str:
    mime = file.mime_type

    # Google native files — already exported as plain text/csv
    if mime in [
        "application/vnd.google-apps.document",
        "application/vnd.google-apps.spreadsheet",
        "application/vnd.google-apps.presentation",
        "text/plain"
    ]:
        return content.decode("utf-8", errors="ignore")

    if mime == "application/pdf":
        return _extract_pdf(content)

    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx(content)

    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _extract_pptx(content)

    return ""


def _extract_pdf(content: bytes) -> str:
    doc = fitz.open(stream=content, filetype="pdf")
    pages = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"[Page {page_num + 1}]\n{text}")
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    doc = docx.Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])


def _extract_pptx(content: bytes) -> str:
    prs = pptx.Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)
